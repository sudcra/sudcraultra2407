from pptx import Presentation

def creappt(nombreppt,nuevoppt,evaluacion, seccion, profesor, itemes):
    # Abre la presentación
    presentation = Presentation("C:/Users/lgutierrez/OneDrive - Fundacion Instituto Profesional Duoc UC/SUDCRA/pptsTemplates/" + nombreppt)

    # Editar el primer slide
    first_slide = presentation.slides[0]
    for shape in first_slide.shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    if "RETRO" in run.text:
                        run.text = "Retroalimentación de la evaluación " + evaluacion
                    if "SECCION" in run.text:
                        run.text = seccion
                    if "profesor" in run.text:
                        run.text = "Profesor: " + profesor
    

    # Eliminar slides 4, 8 y 10
    slides_to_delete = itemes  # Python usa indexado base 0
    for i in range(60,0,-1):
        if i not in itemes:
            try:
                del presentation.slides._sldIdLst[i]
            except:
                    u=1
            
    # Guardar la presentación modificada
    presentation.save("C:/Users/lgutierrez/OneDrive - Fundacion Instituto Profesional Duoc UC/SUDCRA/informes/ppts/" + nuevoppt)

if __name__ == "__main__":

    itemes=[3,5,7,8]
    creappt("MAT3110-2024001-3.pptx", "nuevo8.pptx","Prueba 1", "MAT1111-005D", "Ronny Godoy", itemes)